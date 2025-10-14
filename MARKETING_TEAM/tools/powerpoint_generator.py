"""
PowerPoint Generation Tool
Creates professional presentations and pitch decks
"""

from claude_agent_sdk import tool
from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.enum.text import PP_ALIGN
from pptx.dml.color import RGBColor
import json
import os
from pathlib import Path
from datetime import datetime


@tool(
    "generate_powerpoint",
    "Generate professional PowerPoint presentation",
    {
        "title": str,
        "subtitle": str,
        "slides": list,  # List of slide dicts: [{"type": "title", "title": "...", "content": [...]}]
        "author": str,
        "company": str,
        "filename": str
    }
)
async def generate_powerpoint(args):
    """
    Generate PowerPoint presentation
    """

    title = args["title"]
    subtitle = args.get("subtitle", "")
    slides_data = args["slides"]
    author = args.get("author", "")
    company = args.get("company", "")
    filename = args.get("filename", "presentation.pptx")

    # Ensure output directory exists
    output_dir = Path("output/presentations")
    output_dir.mkdir(parents=True, exist_ok=True)

    # Full output path
    output_path = output_dir / filename

    try:
        # Create presentation
        prs = Presentation()
        prs.slide_width = Inches(16)
        prs.slide_height = Inches(9)

        # Brand colors
        BRAND_PRIMARY = RGBColor(0, 102, 204)  # #0066CC
        BRAND_SECONDARY = RGBColor(0, 204, 153)  # #00CC99
        BRAND_ACCENT = RGBColor(255, 107, 53)  # #FF6B35
        GRAY_DARK = RGBColor(51, 51, 51)
        GRAY_LIGHT = RGBColor(102, 102, 102)

        # --- TITLE SLIDE ---
        title_slide_layout = prs.slide_layouts[0]  # Title slide layout
        slide = prs.slides.add_slide(title_slide_layout)

        # Set title
        title_shape = slide.shapes.title
        title_shape.text = title
        title_shape.text_frame.paragraphs[0].font.size = Pt(54)
        title_shape.text_frame.paragraphs[0].font.bold = True
        title_shape.text_frame.paragraphs[0].font.color.rgb = BRAND_PRIMARY

        # Set subtitle
        if subtitle:
            subtitle_shape = slide.placeholders[1]
            subtitle_shape.text = subtitle
            subtitle_shape.text_frame.paragraphs[0].font.size = Pt(28)
            subtitle_shape.text_frame.paragraphs[0].font.color.rgb = GRAY_LIGHT

        # Add author/company info at bottom
        if author or company:
            info_text = ""
            if author:
                info_text += author
            if company:
                if author:
                    info_text += f" | {company}"
                else:
                    info_text += company

            text_box = slide.shapes.add_textbox(
                Inches(1), Inches(8), Inches(14), Inches(0.5)
            )
            text_frame = text_box.text_frame
            p = text_frame.paragraphs[0]
            p.text = info_text
            p.font.size = Pt(16)
            p.font.color.rgb = GRAY_LIGHT
            p.alignment = PP_ALIGN.CENTER

        # --- CONTENT SLIDES ---
        for slide_data in slides_data:
            slide_type = slide_data.get("type", "content")

            if slide_type == "title_only":
                slide = _create_title_only_slide(prs, slide_data, BRAND_PRIMARY)

            elif slide_type == "content":
                slide = _create_content_slide(prs, slide_data, BRAND_PRIMARY, GRAY_DARK)

            elif slide_type == "two_column":
                slide = _create_two_column_slide(prs, slide_data, BRAND_PRIMARY, GRAY_DARK)

            elif slide_type == "image":
                slide = _create_image_slide(prs, slide_data, BRAND_PRIMARY)

            elif slide_type == "quote":
                slide = _create_quote_slide(prs, slide_data, BRAND_PRIMARY, GRAY_LIGHT)

            elif slide_type == "bullets":
                slide = _create_bullets_slide(prs, slide_data, BRAND_PRIMARY, GRAY_DARK)

            elif slide_type == "closing":
                slide = _create_closing_slide(prs, slide_data, BRAND_PRIMARY, GRAY_LIGHT)

        # Save presentation
        prs.save(str(output_path))

        result = {
            "status": "success",
            "title": title,
            "filename": filename,
            "output_path": str(output_path),
            "file_size": f"{os.path.getsize(output_path) / 1024:.1f} KB",
            "slide_count": len(prs.slides),
            "created_at": datetime.now().isoformat(),
            "next_step": "Use upload_to_google_drive tool to upload and get shareable link"
        }

        return {
            "content": [{
                "type": "text",
                "text": json.dumps(result, indent=2)
            }]
        }

    except Exception as e:
        return {
            "content": [{
                "type": "text",
                "text": f"Error generating PowerPoint: {str(e)}"
            }]
        }


def _create_title_only_slide(prs, slide_data, brand_color):
    """Create title-only slide (section divider)"""
    blank_layout = prs.slide_layouts[6]  # Blank layout
    slide = prs.slides.add_slide(blank_layout)

    # Add centered title
    title_text = slide_data.get("title", "")
    text_box = slide.shapes.add_textbox(
        Inches(2), Inches(3.5), Inches(12), Inches(2)
    )
    text_frame = text_box.text_frame
    p = text_frame.paragraphs[0]
    p.text = title_text
    p.font.size = Pt(48)
    p.font.bold = True
    p.font.color.rgb = brand_color
    p.alignment = PP_ALIGN.CENTER

    return slide


def _create_content_slide(prs, slide_data, brand_color, text_color):
    """Create standard content slide with title and body"""
    title_content_layout = prs.slide_layouts[1]  # Title and content
    slide = prs.slides.add_slide(title_content_layout)

    # Set title
    title = slide.shapes.title
    title.text = slide_data.get("title", "")
    title.text_frame.paragraphs[0].font.size = Pt(40)
    title.text_frame.paragraphs[0].font.bold = True
    title.text_frame.paragraphs[0].font.color.rgb = brand_color

    # Set content
    content_items = slide_data.get("content", [])
    if content_items:
        content_box = slide.placeholders[1]
        text_frame = content_box.text_frame
        text_frame.clear()

        for item in content_items:
            p = text_frame.add_paragraph()
            p.text = item
            p.level = 0
            p.font.size = Pt(24)
            p.font.color.rgb = text_color

    return slide


def _create_two_column_slide(prs, slide_data, brand_color, text_color):
    """Create two-column layout slide"""
    blank_layout = prs.slide_layouts[6]
    slide = prs.slides.add_slide(blank_layout)

    # Add title
    title_box = slide.shapes.add_textbox(Inches(0.5), Inches(0.5), Inches(15), Inches(1))
    text_frame = title_box.text_frame
    p = text_frame.paragraphs[0]
    p.text = slide_data.get("title", "")
    p.font.size = Pt(40)
    p.font.bold = True
    p.font.color.rgb = brand_color

    # Left column
    left_content = slide_data.get("left_content", [])
    left_box = slide.shapes.add_textbox(Inches(0.5), Inches(2), Inches(7), Inches(6))
    text_frame = left_box.text_frame
    for item in left_content:
        p = text_frame.add_paragraph()
        p.text = f"• {item}"
        p.font.size = Pt(20)
        p.font.color.rgb = text_color

    # Right column
    right_content = slide_data.get("right_content", [])
    right_box = slide.shapes.add_textbox(Inches(8.5), Inches(2), Inches(7), Inches(6))
    text_frame = right_box.text_frame
    for item in right_content:
        p = text_frame.add_paragraph()
        p.text = f"• {item}"
        p.font.size = Pt(20)
        p.font.color.rgb = text_color

    return slide


def _create_image_slide(prs, slide_data, brand_color):
    """Create slide with large image and caption"""
    blank_layout = prs.slide_layouts[6]
    slide = prs.slides.add_slide(blank_layout)

    # Add title
    title_text = slide_data.get("title", "")
    if title_text:
        title_box = slide.shapes.add_textbox(Inches(0.5), Inches(0.5), Inches(15), Inches(1))
        text_frame = title_box.text_frame
        p = text_frame.paragraphs[0]
        p.text = title_text
        p.font.size = Pt(40)
        p.font.bold = True
        p.font.color.rgb = brand_color

    # Add image (if path provided)
    image_path = slide_data.get("image_path", "")
    if image_path and os.path.exists(image_path):
        slide.shapes.add_picture(
            image_path,
            Inches(2), Inches(2),
            width=Inches(12)
        )

    # Add caption
    caption = slide_data.get("caption", "")
    if caption:
        caption_box = slide.shapes.add_textbox(Inches(2), Inches(7.5), Inches(12), Inches(0.8))
        text_frame = caption_box.text_frame
        p = text_frame.paragraphs[0]
        p.text = caption
        p.font.size = Pt(18)
        p.alignment = PP_ALIGN.CENTER

    return slide


def _create_quote_slide(prs, slide_data, brand_color, gray_color):
    """Create quote slide"""
    blank_layout = prs.slide_layouts[6]
    slide = prs.slides.add_slide(blank_layout)

    # Add quote
    quote_text = slide_data.get("quote", "")
    quote_box = slide.shapes.add_textbox(Inches(2), Inches(3), Inches(12), Inches(3))
    text_frame = quote_box.text_frame
    p = text_frame.paragraphs[0]
    p.text = f'"{quote_text}"'
    p.font.size = Pt(36)
    p.font.italic = True
    p.font.color.rgb = brand_color
    p.alignment = PP_ALIGN.CENTER

    # Add attribution
    attribution = slide_data.get("attribution", "")
    if attribution:
        attr_box = slide.shapes.add_textbox(Inches(2), Inches(6.5), Inches(12), Inches(1))
        text_frame = attr_box.text_frame
        p = text_frame.paragraphs[0]
        p.text = f"— {attribution}"
        p.font.size = Pt(24)
        p.font.color.rgb = gray_color
        p.alignment = PP_ALIGN.CENTER

    return slide


def _create_bullets_slide(prs, slide_data, brand_color, text_color):
    """Create bullet point slide"""
    title_content_layout = prs.slide_layouts[1]
    slide = prs.slides.add_slide(title_content_layout)

    # Set title
    title = slide.shapes.title
    title.text = slide_data.get("title", "")
    title.text_frame.paragraphs[0].font.size = Pt(40)
    title.text_frame.paragraphs[0].font.bold = True
    title.text_frame.paragraphs[0].font.color.rgb = brand_color

    # Add bullets
    bullets = slide_data.get("bullets", [])
    if bullets:
        content_box = slide.placeholders[1]
        text_frame = content_box.text_frame
        text_frame.clear()

        for bullet in bullets:
            p = text_frame.add_paragraph()
            p.text = bullet
            p.level = 0
            p.font.size = Pt(28)
            p.font.color.rgb = text_color

    return slide


def _create_closing_slide(prs, slide_data, brand_color, gray_color):
    """Create closing/CTA slide"""
    blank_layout = prs.slide_layouts[6]
    slide = prs.slides.add_slide(blank_layout)

    # Main message
    main_text = slide_data.get("main_text", "Thank You")
    main_box = slide.shapes.add_textbox(Inches(2), Inches(3), Inches(12), Inches(2))
    text_frame = main_box.text_frame
    p = text_frame.paragraphs[0]
    p.text = main_text
    p.font.size = Pt(54)
    p.font.bold = True
    p.font.color.rgb = brand_color
    p.alignment = PP_ALIGN.CENTER

    # CTA/Contact info
    cta_text = slide_data.get("cta", "")
    if cta_text:
        cta_box = slide.shapes.add_textbox(Inches(2), Inches(5.5), Inches(12), Inches(2))
        text_frame = cta_box.text_frame
        p = text_frame.paragraphs[0]
        p.text = cta_text
        p.font.size = Pt(28)
        p.font.color.rgb = gray_color
        p.alignment = PP_ALIGN.CENTER

    return slide


@tool(
    "create_pitch_deck",
    "Create a standard 10-slide pitch deck",
    {
        "company_name": str,
        "tagline": str,
        "problem": str,
        "solution": str,
        "market_size": str,
        "traction": str,
        "team": str,
        "ask": str,
        "filename": str
    }
)
async def create_pitch_deck(args):
    """
    Create a Guy Kawasaki-style 10-slide pitch deck
    """
    company = args["company_name"]
    tagline = args["tagline"]

    slides = [
        {
            "type": "content",
            "title": "The Problem",
            "content": [args["problem"]]
        },
        {
            "type": "content",
            "title": "Our Solution",
            "content": [args["solution"]]
        },
        {
            "type": "content",
            "title": "Market Opportunity",
            "content": [args["market_size"]]
        },
        {
            "type": "content",
            "title": "Traction",
            "content": [args["traction"]]
        },
        {
            "type": "content",
            "title": "Team",
            "content": [args["team"]]
        },
        {
            "type": "closing",
            "main_text": "The Ask",
            "cta": args["ask"]
        }
    ]

    return await generate_powerpoint({
        "title": company,
        "subtitle": tagline,
        "slides": slides,
        "company": company,
        "filename": args.get("filename", "pitch-deck.pptx")
    })
