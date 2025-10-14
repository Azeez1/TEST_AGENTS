"""
Presentation Creation Tool using python-pptx
Creates professional PowerPoint presentations programmatically
"""

import os
from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.enum.text import PP_ALIGN
from pptx.dml.color import RGBColor
from PIL import Image


def create_presentation(title, subtitle="", slides_data=None, output_dir="outputs/presentations"):
    """
    Create a professional PowerPoint presentation

    Args:
        title (str): Presentation title
        subtitle (str): Subtitle for title slide
        slides_data (list): List of slide dictionaries with different types
        output_dir (str): Directory to save presentation

    Returns:
        str: Path to created .pptx file

    Slide Types:
        - {"type": "title", "title": "...", "subtitle": "..."}
        - {"type": "content", "title": "...", "bullets": ["...", "..."]}
        - {"type": "content_with_image", "title": "...", "bullets": [...], "image": "path/to/image.png"}
        - {"type": "full_image", "title": "...", "image": "path/to/image.png"}
        - {"type": "two_column", "title": "...", "left": {...}, "right": {...}}
        - {"type": "quote", "quote": "...", "author": "..."}
        - {"type": "chart", "title": "...", "chart_image": "path/to/chart.png"}
    """

    # Ensure output directory exists
    os.makedirs(output_dir, exist_ok=True)

    # Create presentation (16:9 aspect ratio)
    prs = Presentation()
    prs.slide_width = Inches(13.333)  # 16:9 ratio
    prs.slide_height = Inches(7.5)

    # Add title slide
    _add_title_slide(prs, title, subtitle)

    # Add content slides
    if slides_data:
        for slide_data in slides_data:
            slide_type = slide_data.get("type", "content")

            if slide_type == "title":
                _add_title_slide(prs, slide_data.get("title", ""), slide_data.get("subtitle", ""))

            elif slide_type == "content":
                _add_content_slide(prs, slide_data.get("title", ""), slide_data.get("bullets", []))

            elif slide_type == "content_with_image":
                _add_content_with_image_slide(
                    prs,
                    slide_data.get("title", ""),
                    slide_data.get("bullets", []),
                    slide_data.get("image", ""),
                    slide_data.get("image_position", "right")
                )

            elif slide_type == "full_image":
                _add_full_image_slide(
                    prs,
                    slide_data.get("title", ""),
                    slide_data.get("image", "")
                )

            elif slide_type == "two_column":
                _add_two_column_slide(
                    prs,
                    slide_data.get("title", ""),
                    slide_data.get("left", {}),
                    slide_data.get("right", {})
                )

            elif slide_type == "quote":
                _add_quote_slide(
                    prs,
                    slide_data.get("quote", ""),
                    slide_data.get("author", "")
                )

            elif slide_type == "chart":
                _add_chart_slide(
                    prs,
                    slide_data.get("title", ""),
                    slide_data.get("chart_image", "")
                )

    # Save presentation
    safe_filename = "".join(c for c in title[:50] if c.isalnum() or c in (' ', '-', '_')).rstrip()
    safe_filename = safe_filename.replace(' ', '_')
    output_path = os.path.join(output_dir, f"{safe_filename}.pptx")

    prs.save(output_path)
    print(f"[OK] Presentation created: {output_path}")
    return output_path


def _add_title_slide(prs, title, subtitle):
    """Add title slide"""
    slide_layout = prs.slide_layouts[0]  # Title slide layout
    slide = prs.slides.add_slide(slide_layout)

    title_shape = slide.shapes.title
    subtitle_shape = slide.placeholders[1]

    title_shape.text = title
    subtitle_shape.text = subtitle

    # Style title
    title_shape.text_frame.paragraphs[0].font.size = Pt(54)
    title_shape.text_frame.paragraphs[0].font.bold = True
    title_shape.text_frame.paragraphs[0].font.color.rgb = RGBColor(0, 32, 96)  # Dark blue


def _add_content_slide(prs, title, bullets):
    """Add content slide with bullet points"""
    slide_layout = prs.slide_layouts[1]  # Title and content layout
    slide = prs.slides.add_slide(slide_layout)

    title_shape = slide.shapes.title
    title_shape.text = title

    # Add bullets
    content_shape = slide.placeholders[1]
    text_frame = content_shape.text_frame
    text_frame.clear()

    for bullet in bullets:
        p = text_frame.add_paragraph()
        p.text = bullet
        p.level = 0
        p.font.size = Pt(24)
        p.space_before = Pt(12)


def _add_content_with_image_slide(prs, title, bullets, image_path, image_position="right"):
    """Add content slide with bullet points and image side-by-side"""
    slide_layout = prs.slide_layouts[5]  # Blank layout
    slide = prs.slides.add_slide(slide_layout)

    # Add title
    title_box = slide.shapes.add_textbox(Inches(0.5), Inches(0.5), Inches(12), Inches(0.8))
    title_frame = title_box.text_frame
    title_frame.text = title
    title_frame.paragraphs[0].font.size = Pt(44)
    title_frame.paragraphs[0].font.bold = True
    title_frame.paragraphs[0].font.color.rgb = RGBColor(0, 32, 96)

    if image_position == "right":
        # Bullets on left, image on right
        text_left = Inches(0.5)
        text_width = Inches(6)
        image_left = Inches(7)
        image_width = Inches(5.5)
    else:
        # Image on left, bullets on right
        image_left = Inches(0.5)
        image_width = Inches(5.5)
        text_left = Inches(6.5)
        text_width = Inches(6)

    # Add bullets
    text_box = slide.shapes.add_textbox(text_left, Inches(1.5), text_width, Inches(5))
    text_frame = text_box.text_frame
    for bullet in bullets:
        p = text_frame.add_paragraph()
        p.text = bullet
        p.level = 0
        p.font.size = Pt(20)
        p.space_before = Pt(10)

    # Add image
    if image_path and os.path.exists(image_path):
        slide.shapes.add_picture(image_path, image_left, Inches(1.5), width=image_width)


def _add_full_image_slide(prs, title, image_path):
    """Add slide with full-screen background image and title overlay"""
    slide_layout = prs.slide_layouts[6]  # Blank layout
    slide = prs.slides.add_slide(slide_layout)

    # Add image as background
    if image_path and os.path.exists(image_path):
        slide.shapes.add_picture(
            image_path,
            Inches(0), Inches(0),
            width=prs.slide_width,
            height=prs.slide_height
        )

    # Add title overlay with semi-transparent background
    title_box = slide.shapes.add_textbox(Inches(0.5), Inches(3), Inches(12), Inches(1.5))
    title_frame = title_box.text_frame
    title_frame.text = title
    title_frame.paragraphs[0].font.size = Pt(54)
    title_frame.paragraphs[0].font.bold = True
    title_frame.paragraphs[0].font.color.rgb = RGBColor(255, 255, 255)
    title_frame.paragraphs[0].alignment = PP_ALIGN.CENTER


def _add_two_column_slide(prs, title, left_content, right_content):
    """Add slide with two columns"""
    slide_layout = prs.slide_layouts[5]  # Blank layout
    slide = prs.slides.add_slide(slide_layout)

    # Add title
    title_box = slide.shapes.add_textbox(Inches(0.5), Inches(0.5), Inches(12), Inches(0.8))
    title_frame = title_box.text_frame
    title_frame.text = title
    title_frame.paragraphs[0].font.size = Pt(44)
    title_frame.paragraphs[0].font.bold = True

    # Left column
    left_box = slide.shapes.add_textbox(Inches(0.5), Inches(1.5), Inches(5.5), Inches(5))
    left_frame = left_box.text_frame
    left_frame.text = left_content.get("text", "")
    if left_content.get("image") and os.path.exists(left_content["image"]):
        slide.shapes.add_picture(left_content["image"], Inches(0.5), Inches(2.5), width=Inches(5.5))

    # Right column
    right_box = slide.shapes.add_textbox(Inches(7), Inches(1.5), Inches(5.5), Inches(5))
    right_frame = right_box.text_frame
    right_frame.text = right_content.get("text", "")
    if right_content.get("image") and os.path.exists(right_content["image"]):
        slide.shapes.add_picture(right_content["image"], Inches(7), Inches(2.5), width=Inches(5.5))


def _add_quote_slide(prs, quote, author):
    """Add quote/testimonial slide"""
    slide_layout = prs.slide_layouts[6]  # Blank layout
    slide = prs.slides.add_slide(slide_layout)

    # Add quote
    quote_box = slide.shapes.add_textbox(Inches(2), Inches(2.5), Inches(9), Inches(2))
    quote_frame = quote_box.text_frame
    quote_frame.text = f'"{quote}"'
    quote_frame.paragraphs[0].font.size = Pt(36)
    quote_frame.paragraphs[0].font.italic = True
    quote_frame.paragraphs[0].alignment = PP_ALIGN.CENTER

    # Add author
    author_box = slide.shapes.add_textbox(Inches(2), Inches(5), Inches(9), Inches(0.8))
    author_frame = author_box.text_frame
    author_frame.text = f"— {author}"
    author_frame.paragraphs[0].font.size = Pt(24)
    author_frame.paragraphs[0].alignment = PP_ALIGN.CENTER


def _add_chart_slide(prs, title, chart_image_path):
    """Add slide with chart/data visualization"""
    slide_layout = prs.slide_layouts[5]  # Blank layout
    slide = prs.slides.add_slide(slide_layout)

    # Add title
    title_box = slide.shapes.add_textbox(Inches(0.5), Inches(0.5), Inches(12), Inches(0.8))
    title_frame = title_box.text_frame
    title_frame.text = title
    title_frame.paragraphs[0].font.size = Pt(44)
    title_frame.paragraphs[0].font.bold = True

    # Add chart image
    if chart_image_path and os.path.exists(chart_image_path):
        slide.shapes.add_picture(
            chart_image_path,
            Inches(1.5), Inches(1.5),
            width=Inches(10)
        )


# Test function
if __name__ == "__main__":
    # Test presentation creation
    test_slides = [
        {
            "type": "content",
            "title": "Key Objectives",
            "bullets": [
                "Increase brand awareness by 35%",
                "Drive 25% growth in qualified leads",
                "Expand market share in key demographics"
            ]
        },
        {
            "type": "content",
            "title": "Campaign Overview",
            "bullets": [
                "Multi-channel campaign launch",
                "Account-based marketing for enterprise clients",
                "Influencer partnerships"
            ]
        }
    ]

    presentation_path = create_presentation(
        title="Test Marketing Presentation",
        subtitle="Demo Presentation",
        slides_data=test_slides
    )

    print(f"Test presentation created: {presentation_path}")
