"""
Create AI Basics PowerPoint with Botanical Garden Theme
Uses python-pptx to create a 5-slide presentation about AI
"""

from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.enum.text import PP_ALIGN
from pptx.dml.color import RGBColor

# Botanical Garden Theme Colors
FERN_GREEN = RGBColor(74, 124, 89)      # #4a7c59 - Rich natural green
MARIGOLD = RGBColor(249, 166, 32)       # #f9a620 - Bright floral accent
TERRACOTTA = RGBColor(183, 71, 42)      # #b7472a - Earthy warm tone
CREAM = RGBColor(245, 243, 237)         # #f5f3ed - Soft neutral backgrounds

def set_background_color(slide, color):
    """Set slide background to solid color"""
    background = slide.background
    fill = background.fill
    fill.solid()
    fill.fore_color.rgb = color

def add_title_slide(prs):
    """Slide 1: Title slide"""
    slide = prs.slides.add_slide(prs.slide_layouts[6])  # Blank layout
    set_background_color(slide, CREAM)

    # Title
    title_box = slide.shapes.add_textbox(Inches(1), Inches(2.5), Inches(8), Inches(1.5))
    title_frame = title_box.text_frame
    title_frame.text = "Artificial Intelligence:\nGrowing the Future"
    title_para = title_frame.paragraphs[0]
    title_para.alignment = PP_ALIGN.CENTER
    title_para.font.size = Pt(44)
    title_para.font.bold = True
    title_para.font.color.rgb = FERN_GREEN

    # Subtitle
    subtitle_box = slide.shapes.add_textbox(Inches(2), Inches(4.2), Inches(6), Inches(0.8))
    subtitle_frame = subtitle_box.text_frame
    subtitle_frame.text = "Cultivating Intelligence Through Technology"
    subtitle_para = subtitle_frame.paragraphs[0]
    subtitle_para.alignment = PP_ALIGN.CENTER
    subtitle_para.font.size = Pt(24)
    subtitle_para.font.color.rgb = TERRACOTTA

    # Decorative accent bar
    accent_shape = slide.shapes.add_shape(
        1,  # Rectangle
        Inches(3.5), Inches(5.2), Inches(3), Inches(0.1)
    )
    accent_shape.fill.solid()
    accent_shape.fill.fore_color.rgb = MARIGOLD
    accent_shape.line.fill.background()

def add_what_is_ai_slide(prs):
    """Slide 2: What is AI?"""
    slide = prs.slides.add_slide(prs.slide_layouts[6])  # Blank layout
    set_background_color(slide, CREAM)

    # Header bar
    header_shape = slide.shapes.add_shape(
        1,  # Rectangle
        Inches(0), Inches(0), Inches(10), Inches(1)
    )
    header_shape.fill.solid()
    header_shape.fill.fore_color.rgb = FERN_GREEN
    header_shape.line.fill.background()

    # Title
    title_box = slide.shapes.add_textbox(Inches(0.5), Inches(0.25), Inches(9), Inches(0.5))
    title_frame = title_box.text_frame
    title_frame.text = "What is AI?"
    title_para = title_frame.paragraphs[0]
    title_para.font.size = Pt(36)
    title_para.font.bold = True
    title_para.font.color.rgb = CREAM

    # Content
    content_box = slide.shapes.add_textbox(Inches(1), Inches(1.5), Inches(8), Inches(4.5))
    content_frame = content_box.text_frame
    content_frame.word_wrap = True

    # Definition
    p1 = content_frame.paragraphs[0]
    p1.text = "AI is technology that enables machines to learn and make decisions"
    p1.font.size = Pt(20)
    p1.font.color.rgb = FERN_GREEN
    p1.space_after = Pt(20)

    # Subtitle
    p2 = content_frame.add_paragraph()
    p2.text = "Like a garden, AI grows through:"
    p2.font.size = Pt(18)
    p2.font.bold = True
    p2.font.color.rgb = TERRACOTTA
    p2.space_after = Pt(12)

    # Bullet points
    bullets = [
        "Learning from data (seeds of knowledge)",
        "Recognizing patterns (seasonal cycles)",
        "Making predictions (forecasting growth)",
        "Improving over time (natural evolution)"
    ]

    for bullet in bullets:
        p = content_frame.add_paragraph()
        p.text = bullet
        p.level = 0
        p.font.size = Pt(18)
        p.font.color.rgb = FERN_GREEN
        p.space_after = Pt(8)

def add_how_ai_works_slide(prs):
    """Slide 3: How AI Works"""
    slide = prs.slides.add_slide(prs.slide_layouts[6])  # Blank layout
    set_background_color(slide, CREAM)

    # Header bar
    header_shape = slide.shapes.add_shape(
        1,  # Rectangle
        Inches(0), Inches(0), Inches(10), Inches(1)
    )
    header_shape.fill.solid()
    header_shape.fill.fore_color.rgb = FERN_GREEN
    header_shape.line.fill.background()

    # Title
    title_box = slide.shapes.add_textbox(Inches(0.5), Inches(0.25), Inches(9), Inches(0.5))
    title_frame = title_box.text_frame
    title_frame.text = "How AI Works: The Growth Process"
    title_para = title_frame.paragraphs[0]
    title_para.font.size = Pt(36)
    title_para.font.bold = True
    title_para.font.color.rgb = CREAM

    # Content
    content_box = slide.shapes.add_textbox(Inches(1), Inches(1.5), Inches(8), Inches(4.5))
    content_frame = content_box.text_frame
    content_frame.word_wrap = True

    # Three main concepts
    concepts = [
        ("1. Data (The Soil)", [
            "AI needs quality data to learn",
            "More data = better understanding",
            "Clean data produces better results"
        ]),
        ("2. Learning (The Roots)", [
            "Algorithms find patterns in data",
            "System adapts and improves",
            "Deeper connections form over time"
        ]),
        ("3. Predictions (The Harvest)", [
            "AI makes informed decisions",
            "Provides recommendations",
            "Solves complex problems"
        ])
    ]

    for i, (heading, points) in enumerate(concepts):
        if i > 0:
            content_frame.add_paragraph()  # Spacing

        # Heading
        p_heading = content_frame.add_paragraph() if i > 0 else content_frame.paragraphs[0]
        p_heading.text = heading
        p_heading.font.size = Pt(20)
        p_heading.font.bold = True
        p_heading.font.color.rgb = MARIGOLD
        p_heading.space_after = Pt(6)

        # Points
        for point in points:
            p = content_frame.add_paragraph()
            p.text = point
            p.level = 1
            p.font.size = Pt(16)
            p.font.color.rgb = FERN_GREEN
            p.space_after = Pt(4)

def add_ai_in_daily_life_slide(prs):
    """Slide 4: AI in Daily Life"""
    slide = prs.slides.add_slide(prs.slide_layouts[6])  # Blank layout
    set_background_color(slide, CREAM)

    # Header bar
    header_shape = slide.shapes.add_shape(
        1,  # Rectangle
        Inches(0), Inches(0), Inches(10), Inches(1)
    )
    header_shape.fill.solid()
    header_shape.fill.fore_color.rgb = FERN_GREEN
    header_shape.line.fill.background()

    # Title
    title_box = slide.shapes.add_textbox(Inches(0.5), Inches(0.25), Inches(9), Inches(0.5))
    title_frame = title_box.text_frame
    title_frame.text = "AI in Daily Life: Where It Blooms"
    title_para = title_frame.paragraphs[0]
    title_para.font.size = Pt(36)
    title_para.font.bold = True
    title_para.font.color.rgb = CREAM

    # Content
    content_box = slide.shapes.add_textbox(Inches(1), Inches(1.5), Inches(8), Inches(4.5))
    content_frame = content_box.text_frame
    content_frame.word_wrap = True

    # Applications
    applications = [
        ("Voice Assistants", "Siri, Alexa, Google Assistant help with daily tasks"),
        ("Recommendations", "Netflix, Spotify, Amazon suggest content you'll enjoy"),
        ("Smart Photos", "Your phone organizes and enhances photos automatically"),
        ("Navigation", "Maps predict traffic and find optimal routes"),
        ("Email Filters", "Spam detection keeps your inbox clean"),
        ("Health Monitoring", "Fitness trackers analyze activity and wellness"),
        ("Shopping", "Chatbots provide customer support 24/7"),
        ("Social Media", "Feed curation shows relevant content")
    ]

    p = content_frame.paragraphs[0]
    p.text = "AI is already part of your everyday life:"
    p.font.size = Pt(18)
    p.font.bold = True
    p.font.color.rgb = TERRACOTTA
    p.space_after = Pt(12)

    for app, description in applications:
        p_app = content_frame.add_paragraph()
        p_app.text = f"{app}: "
        p_app.font.size = Pt(16)
        p_app.font.bold = True
        p_app.font.color.rgb = MARIGOLD
        p_app.space_after = Pt(2)

        # Add description as run
        run = p_app.add_run()
        run.text = description
        run.font.bold = False
        run.font.color.rgb = FERN_GREEN

def add_future_slide(prs):
    """Slide 5: The Future of AI"""
    slide = prs.slides.add_slide(prs.slide_layouts[6])  # Blank layout
    set_background_color(slide, CREAM)

    # Header bar
    header_shape = slide.shapes.add_shape(
        1,  # Rectangle
        Inches(0), Inches(0), Inches(10), Inches(1)
    )
    header_shape.fill.solid()
    header_shape.fill.fore_color.rgb = FERN_GREEN
    header_shape.line.fill.background()

    # Title
    title_box = slide.shapes.add_textbox(Inches(0.5), Inches(0.25), Inches(9), Inches(0.5))
    title_frame = title_box.text_frame
    title_frame.text = "The Future of AI: Endless Growth"
    title_para = title_frame.paragraphs[0]
    title_para.font.size = Pt(36)
    title_para.font.bold = True
    title_para.font.color.rgb = CREAM

    # Content
    content_box = slide.shapes.add_textbox(Inches(1), Inches(1.5), Inches(8), Inches(4.5))
    content_frame = content_box.text_frame
    content_frame.word_wrap = True

    # Future opportunities
    p1 = content_frame.paragraphs[0]
    p1.text = "Opportunities on the Horizon:"
    p1.font.size = Pt(20)
    p1.font.bold = True
    p1.font.color.rgb = TERRACOTTA
    p1.space_after = Pt(12)

    opportunities = [
        "Healthcare: Early disease detection and personalized medicine",
        "Education: Adaptive learning tailored to each student",
        "Environment: Climate modeling and conservation efforts",
        "Transportation: Self-driving vehicles and smart cities",
        "Creativity: AI as a collaborative creative partner",
        "Accessibility: Breaking down barriers for people with disabilities"
    ]

    for opp in opportunities:
        p = content_frame.add_paragraph()
        p.text = opp
        p.level = 0
        p.font.size = Pt(16)
        p.font.color.rgb = FERN_GREEN
        p.space_after = Pt(6)

    # Closing statement
    content_frame.add_paragraph()
    p_closing = content_frame.add_paragraph()
    p_closing.text = "Like a thriving garden, AI continues to grow and flourish,\ncreating new possibilities for humanity."
    p_closing.font.size = Pt(16)
    p_closing.font.italic = True
    p_closing.font.color.rgb = MARIGOLD
    p_closing.alignment = PP_ALIGN.CENTER
    p_closing.space_before = Pt(12)

def main():
    """Create the presentation"""
    prs = Presentation()
    prs.slide_width = Inches(10)
    prs.slide_height = Inches(7.5)

    # Add all slides
    add_title_slide(prs)
    add_what_is_ai_slide(prs)
    add_how_ai_works_slide(prs)
    add_ai_in_daily_life_slide(prs)
    add_future_slide(prs)

    # Save presentation
    output_path = r"c:\Users\sabaa\OneDrive\Desktop\TEST_AGENTS\MARKETING_TEAM\outputs\presentations\ai_basics_botanical.pptx"
    prs.save(output_path)
    print(f"✅ Presentation created successfully!")
    print(f"📁 Location: {output_path}")
    print(f"🎨 Theme: Botanical Garden (Fern Green, Marigold, Terracotta, Cream)")
    print(f"📊 Slides: 5")

if __name__ == "__main__":
    main()
