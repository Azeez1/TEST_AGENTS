"""Rebuild NotebookLM PPTX from scratch with watermark cropped off."""
import os, glob, re
from PIL import Image
from pptx import Presentation
from pptx.util import Emu

SRC_DIR = r'C:\Users\sabaa\ONEDRIVE\DESKTOP\TEST_AGENTS\tmp_slides'
OUT_DIR = r'C:\Users\sabaa\ONEDRIVE\DESKTOP\TEST_AGENTS\tmp_slides_cleaned'
OUT_PPTX = r'C:\Users\sabaa\ONEDRIVE\DESKTOP\TEST_AGENTS\MARKETING_TEAM\outputs\social_media\natural_body_teardown_carousel.pptx'

os.makedirs(OUT_DIR, exist_ok=True)

def nat_key(s):
    m = re.search(r'(\d+)', os.path.basename(s))
    return int(m.group(1)) if m else 0

KEEP = [1, 4, 5, 6, 8, 10, 13, 14]
images = sorted(glob.glob(os.path.join(SRC_DIR, 'image*.png')), key=nat_key)
images = [p for p in images if nat_key(p) in KEEP]
images.sort(key=lambda p: KEEP.index(nat_key(p)))
print(f'Using {len(images)} slides: {[nat_key(p) for p in images]}')

# Crop bottom-right watermark zone: remove bottom ~35px which contains the NotebookLM text
cleaned_paths = []
for img_path in images:
    img = Image.open(img_path)
    w, h = img.size
    # Watermark is in far bottom-right corner. Crop a small rect on bottom-right only
    # by masking with sampled bg color from an area we know is clean
    from PIL import ImageDraw
    # Sample a clean area: middle-right vertical band (safe zone between watermarks)
    sample_box = (w - 80, int(h * 0.4), w - 10, int(h * 0.6))
    sample = img.crop(sample_box)
    sample_small = sample.resize((1, 1))
    bg_color = sample_small.getpixel((0, 0))
    # Mask TOP-RIGHT corner where NotebookLM watermark sits (small box, ~5% x 4%)
    img = img.convert('RGB')
    draw = ImageDraw.Draw(img)
    mask_box = (int(w * 0.90), int(h * 0.92), w, h)
    draw.rectangle(mask_box, fill=bg_color)
    out_path = os.path.join(OUT_DIR, os.path.basename(img_path))
    img.save(out_path, 'PNG')
    cleaned_paths.append(out_path)

print(f'Cleaned {len(cleaned_paths)} images -> {OUT_DIR}')

# Build fresh PPTX with 16:9 widescreen, each slide = full-page image
prs = Presentation()
prs.slide_width = Emu(16256000)
prs.slide_height = Emu(9144000)
blank_layout = prs.slide_layouts[6]  # blank

for img_path in cleaned_paths:
    slide = prs.slides.add_slide(blank_layout)
    slide.shapes.add_picture(img_path, 0, 0, width=prs.slide_width, height=prs.slide_height)

prs.save(OUT_PPTX)
size_mb = os.path.getsize(OUT_PPTX) / 1024 / 1024
print(f'\nBuilt: {OUT_PPTX}')
print(f'Slides: {len(cleaned_paths)}')
print(f'Size: {size_mb:.1f} MB')
